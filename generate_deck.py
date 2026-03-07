"""
Generate the Microsoft Agent Framework Workshop PowerPoint deck.
Run: python generate_deck.py
Output: MAF-Workshop-Deck.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Brand colors ──────────────────────────────────────────────────────────────
MSFT_BLUE = RGBColor(0x00, 0x78, 0xD4)
DARK_BLUE = RGBColor(0x00, 0x2B, 0x5C)
LIGHT_BLUE = RGBColor(0xDE, 0xEC, 0xF9)
ACCENT_GREEN = RGBColor(0x10, 0x7C, 0x10)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
ACCENT_RED = RGBColor(0xD1, 0x34, 0x38)
ACCENT_PURPLE = RGBColor(0x6B, 0x69, 0xD6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MED_GRAY = RGBColor(0x60, 0x60, 0x60)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
CODE_BG = RGBColor(0x1E, 0x1E, 0x1E)
CODE_FG = RGBColor(0xD4, 0xD4, 0xD4)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


# ── Helper functions ──────────────────────────────────────────────────────────

def add_bg_rect(slide, color, left=0, top=0, width=None, height=None):
    """Add a filled rectangle shape as a background element."""
    width = width or W
    height = height or H
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=DARK_GRAY, align=PP_ALIGN.LEFT,
                 font_name="Segoe UI"):
    """Add a simple text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txBox


def add_bullet_slide(slide, items, left, top, width, height,
                     font_size=18, color=DARK_GRAY, spacing=Pt(8),
                     bold_prefix=False):
    """Add a text box with bulleted items."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = spacing
        p.level = 0

        # Support bold prefix with "**prefix** rest" syntax
        if bold_prefix and "—" in item:
            parts = item.split("—", 1)
            run1 = p.add_run()
            run1.text = parts[0].strip()
            run1.font.bold = True
            run1.font.size = Pt(font_size)
            run1.font.color.rgb = color
            run1.font.name = "Segoe UI"
            run2 = p.add_run()
            run2.text = " — " + parts[1].strip()
            run2.font.size = Pt(font_size)
            run2.font.color.rgb = color
            run2.font.name = "Segoe UI"
        else:
            run = p.add_run()
            run.text = item
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.name = "Segoe UI"
    return txBox


def add_code_block(slide, code, left, top, width, height, font_size=13):
    """Add a code block with dark background."""
    bg = add_bg_rect(slide, CODE_BG, left, top, width, height)
    bg.shadow.inherit = False

    txBox = slide.shapes.add_textbox(
        left + Inches(0.2), top + Inches(0.15),
        width - Inches(0.4), height - Inches(0.3))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(code.strip().split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(2)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.color.rgb = CODE_FG
        run.font.name = "Cascadia Code"


def new_slide():
    """Create a blank slide."""
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def section_title_slide(section_num, section_title, time_range, color=MSFT_BLUE):
    """Create a section divider slide."""
    slide = new_slide()
    add_bg_rect(slide, color)
    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(1),
                 section_num, font_size=24, color=WHITE, bold=False)
    add_text_box(slide, Inches(0.8), Inches(2.3), Inches(11), Inches(2),
                 section_title, font_size=44, color=WHITE, bold=True)
    add_text_box(slide, Inches(0.8), Inches(4.5), Inches(11), Inches(0.8),
                 time_range, font_size=20, color=RGBColor(0xBB, 0xDE, 0xFB))
    return slide


def content_slide(title, subtitle=None):
    """Create a standard content slide with title bar."""
    slide = new_slide()
    # Blue top bar
    add_bg_rect(slide, MSFT_BLUE, Inches(0), Inches(0), W, Inches(1.15))
    add_text_box(slide, Inches(0.6), Inches(0.2), Inches(12), Inches(0.75),
                 title, font_size=28, bold=True, color=WHITE)
    if subtitle:
        add_text_box(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(0.5),
                     subtitle, font_size=16, color=MED_GRAY)
    return slide


def table_slide_helper(slide, headers, rows, left, top, col_widths):
    """Add a simple table to a slide."""
    num_rows = len(rows) + 1
    num_cols = len(headers)
    total_w = sum(col_widths)
    tbl_shape = slide.shapes.add_table(num_rows, num_cols, left, top,
                                        Emu(sum(Inches(w).emu for w in col_widths)),
                                        Inches(0.45 * num_rows))
    table = tbl_shape.table

    for ci in range(num_cols):
        table.columns[ci].width = Inches(col_widths[ci])

    # Header row
    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = "Segoe UI"
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE

    # Data rows
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13)
                p.font.color.rgb = DARK_GRAY
                p.font.name = "Segoe UI"
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if ri % 2 == 0 else LIGHT_GRAY

    return tbl_shape


def add_notes(slide, notes_text):
    """Add speaker notes to a slide."""
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes_text


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1: Title Slide
# ══════════════════════════════════════════════════════════════════════════════
slide = new_slide()
add_bg_rect(slide, DARK_BLUE)
# Accent stripe
add_bg_rect(slide, MSFT_BLUE, Inches(0), Inches(3.2), W, Inches(0.08))

add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(1.2),
             "Microsoft Agent Framework", font_size=48, bold=True, color=WHITE)
add_text_box(slide, Inches(0.8), Inches(2.3), Inches(11), Inches(0.8),
             "From Concepts to Code — Full-Day Hands-On Workshop",
             font_size=24, color=RGBColor(0xBB, 0xDE, 0xFB))
add_text_box(slide, Inches(0.8), Inches(4.2), Inches(11), Inches(1.5),
             "[Presenter Names]\n[Partner Name]  ·  [Date]",
             font_size=18, color=RGBColor(0x99, 0xBB, 0xDD))
add_notes(slide, "Welcome everyone. Housekeeping: Wi-Fi, restrooms, breaks schedule. Introduce yourself and co-presenter.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2: Agenda
# ══════════════════════════════════════════════════════════════════════════════
slide = content_slide("Today's Journey")
items = [
    "09:00 – 09:30   Welcome & Intro: Why Agentic AI?",
    "09:30 – 10:15   Module 1 — Agent Fundamentals  [Theory + Lab]",
    "10:15 – 10:30   ☕ Break",
    "10:30 – 11:15   Module 2 — Tools & Function Calling  [Theory + Lab]",
    "11:15 – 12:00   Module 3 — Conversations, Memory & Middleware  [Theory + Lab]",
    "12:00 – 12:45   🍽️ Lunch",
    "12:45 – 13:30   Module 4 — Workflows & Multi-Agent Patterns  [Theory + Lab]",
    "13:30 – 13:45   Hackathon Briefing",
    "14:00 – 16:30   🏗️ Hackathon: Add Agentic AI to the Existing App",
    "16:30 – 17:00   Show & Tell + Wrap-up",
]
add_bullet_slide(slide, items, Inches(0.6), Inches(1.5), Inches(12), Inches(5.5),
                 font_size=19, color=DARK_GRAY)
add_notes(slide, "Set expectations: by noon you'll have built 4 agents, by 5pm you'll have added agentic AI to a real app.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3: AI Agent Evolution
# ══════════════════════════════════════════════════════════════════════════════
slide = content_slide("From Chatbots to Autonomous Agents")
# Three columns showing evolution
for i, (label, desc, clr) in enumerate([
    ("Chatbot", "Rule-based\nScripted responses\nNo reasoning", MED_GRAY),
    ("Copilot", "LLM-powered\nHuman-directed\nSuggests actions", ACCENT_PURPLE),
    ("Agent", "Autonomous\nGoal-driven\nUses tools\nMaintains memory", MSFT_BLUE),
]):
    x = Inches(0.8 + i * 4.2)
    box = add_bg_rect(slide, clr, x, Inches(2.0), Inches(3.6), Inches(4.2))
    box.shadow.inherit = False
    add_text_box(slide, x + Inches(0.3), Inches(2.3), Inches(3.0), Inches(0.7),
                 label, font_size=28, bold=True, color=WHITE)
    add_text_box(slide, x + Inches(0.3), Inches(3.2), Inches(3.0), Inches(2.5),
                 desc, font_size=18, color=WHITE)

# Arrows between columns
for i in range(2):
    x = Inches(4.6 + i * 4.2)
    add_text_box(slide, x, Inches(3.5), Inches(0.6), Inches(0.6),
                 "→", font_size=36, bold=True, color=DARK_GRAY, align=PP_ALIGN.CENTER)

add_notes(slide, "The key shift is autonomy — agents decide WHAT to do, not just HOW to respond. Each generation builds on the previous.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4: Four Properties
# ══════════════════════════════════════════════════════════════════════════════
slide = content_slide("Four Properties of an AI Agent")
props = [
    ("🎯", "Goal-driven", "Works toward a specific objective autonomously"),
    ("🔧", "Tool-using", "Calls functions, APIs, databases, and external services"),
    ("🧠", "Reasoning", "Plans multi-step actions and adapts based on results"),
    ("💾", "Memory", "Maintains context across interactions and sessions"),
]
for i, (icon, title, desc) in enumerate(props):
    y = Inches(1.7 + i * 1.35)
    add_text_box(slide, Inches(0.8), y, Inches(0.8), Inches(0.6),
                 icon, font_size=32, align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.8), y, Inches(3.0), Inches(0.5),
                 title, font_size=22, bold=True, color=DARK_BLUE)
    add_text_box(slide, Inches(1.8), y + Inches(0.45), Inches(10), Inches(0.5),
                 desc, font_size=17, color=MED_GRAY)
add_notes(slide, "Not all AI apps need agents. If a simple API call suffices, use that instead. Agents add value when tasks are complex, multi-step, and require judgment.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5: Use Cases
# ══════════════════════════════════════════════════════════════════════════════
slide = content_slide("Where Agents Shine")
cases = [
    ("💬", "Customer Support", "Resolves tickets autonomously, escalates when needed"),
    ("⚙️", "Process Automation", "Expense approvals, order processing, HR onboarding"),
    ("💻", "Code Assistant", "Writes, reviews, tests, and deploys code"),
    ("🔍", "Enterprise Search", "Finds, summarizes, and synthesizes documents"),
]
for i, (icon, title, desc) in enumerate(cases):
    col = i % 2
    row = i // 2
    x = Inches(0.6 + col * 6.3)
    y = Inches(1.7 + row * 2.5)
    bg = add_bg_rect(slide, LIGHT_BLUE, x, y, Inches(5.8), Inches(2.0))
    bg.shadow.inherit = False
    add_text_box(slide, x + Inches(0.2), y + Inches(0.2), Inches(0.7), Inches(0.6),
                 icon, font_size=32)
    add_text_box(slide, x + Inches(1.0), y + Inches(0.2), Inches(4.5), Inches(0.5),
                 title, font_size=20, bold=True, color=DARK_BLUE)
    add_text_box(slide, x + Inches(1.0), y + Inches(0.8), Inches(4.5), Inches(0.9),
                 desc, font_size=16, color=MED_GRAY)
add_notes(slide, "Relate to the partner's industry if possible. Ask the audience which use case resonates most.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6: Microsoft AI Stack
# ══════════════════════════════════════════════════════════════════════════════
slide = content_slide("Where Does MAF Fit?", "The Microsoft AI Development Stack")
layers = [
    (Inches(1.6), "Copilot Studio", "No-code / low-code agents", RGBColor(0x6B, 0x69, 0xD6)),
    (Inches(3.0), "M365 Agents SDK", "Teams & Microsoft 365 agents", RGBColor(0x00, 0x99, 0xBC)),
    (Inches(4.4), "Microsoft Agent Framework", "Pro-code, full control, open source", MSFT_BLUE),
    (Inches(5.8), "Azure AI Foundry  ·  Azure OpenAI", "Foundation models & infrastructure", DARK_BLUE),
]
for y, title, desc, clr in layers:
    bg = add_bg_rect(slide, clr, Inches(1.5), y, Inches(10.3), Inches(1.15))
    bg.shadow.inherit = False
    add_text_box(slide, Inches(1.8), y + Inches(0.1), Inches(5), Inches(0.5),
                 title, font_size=20, bold=True, color=WHITE)
    add_text_box(slide, Inches(7.0), y + Inches(0.15), Inches(4.5), Inches(0.5),
                 desc, font_size=15, color=RGBColor(0xDD, 0xEE, 0xFF))
add_text_box(slide, Inches(1.5), y + Inches(1.3), Inches(10), Inches(0.5),
             "⬆ You are here: maximum flexibility and control for developers",
             font_size=14, color=ACCENT_ORANGE, bold=True, align=PP_ALIGN.CENTER)
add_notes(slide, "MAF is for developers who need maximum flexibility and control. Copilot Studio is for citizen developers. M365 SDK for Teams-specific scenarios.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7: SK + AutoGen → MAF
# ══════════════════════════════════════════════════════════════════════════════
slide = content_slide("The Best of Both Worlds", "Semantic Kernel + AutoGen → Agent Framework")
# Three columns
cols = [
    ("Semantic Kernel", "Enterprise-grade\nDependency injection\nTelemetry & logging\nType safety\nExtensive model support", RGBColor(0x6B, 0x69, 0xD6)),
    ("AutoGen", "Multi-agent patterns\nSimple abstractions\nResearch innovations\nDynamic collaboration\nOpen-ended orchestration", ACCENT_GREEN),
    ("Agent Framework", "All of the above, plus:\nGraph-based Workflows\nMCP & AG-UI protocols\nSession management\nMiddleware pipeline\nAgent-as-Tool", MSFT_BLUE),
]
for i, (title, items, clr) in enumerate(cols):
    x = Inches(0.5 + i * 4.2)
    bg = add_bg_rect(slide, clr, x, Inches(1.7), Inches(3.9), Inches(5.0))
    bg.shadow.inherit = False
    add_text_box(slide, x + Inches(0.3), Inches(1.9), Inches(3.4), Inches(0.6),
                 title, font_size=22, bold=True, color=WHITE)
    add_text_box(slide, x + Inches(0.3), Inches(2.7), Inches(3.4), Inches(3.8),
                 items, font_size=15, color=WHITE)

# Plus signs
for i in range(2):
    x = Inches(4.5 + i * 4.2)
    label = "+" if i == 0 else "="
    add_text_box(slide, x, Inches(3.5), Inches(0.5), Inches(0.6),
                 label, font_size=36, bold=True, color=DARK_GRAY, align=PP_ALIGN.CENTER)
add_notes(slide, "Same teams, unified vision. Semantic Kernel and AutoGen are predecessors. MAF is the direct successor combining both.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8: Architecture
# ══════════════════════════════════════════════════════════════════════════════
slide = content_slide("MAF Architecture at a Glance")
components = [
    (0, 0, "Agent\n(AIAgent)", MSFT_BLUE),
    (1, 0, "Tools\n(Functions, MCP)", ACCENT_GREEN),
    (2, 0, "Middleware\n(Logging, Security)", ACCENT_PURPLE),
    (0, 1, "Session\n(AgentSession)", RGBColor(0x00, 0x99, 0xBC)),
    (1, 1, "Providers\n(Azure OpenAI, Ollama…)", DARK_BLUE),
    (2, 1, "Workflows\n(Graph orchestration)", ACCENT_ORANGE),
]
for col, row, label, clr in components:
    x = Inches(0.6 + col * 4.2)
    y = Inches(1.7 + row * 2.7)
    bg = add_bg_rect(slide, clr, x, y, Inches(3.8), Inches(2.2))
    bg.shadow.inherit = False
    add_text_box(slide, x + Inches(0.2), y + Inches(0.5), Inches(3.4), Inches(1.5),
                 label, font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_notes(slide, "This diagram will be our compass throughout the day. Keep it visible or printable. Each module goes deep on one or two of these boxes.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9: Open Protocols
# ══════════════════════════════════════════════════════════════════════════════
slide = content_slide("Built on Open Standards")
protocols = [
    ("MCP", "Model Context Protocol", "Standardized tool discovery & invocation.\nAny MCP server works with any MCP client.\nThe universal adapter for agent tools."),
    ("AG-UI", "Agent-to-UI Protocol", "Rich, real-time agent ↔ UI communication.\nStreaming responses, status updates, tool calls.\nBuild interactive web frontends."),
    ("A2A", "Agent-to-Agent Protocol", "Cross-system agent interoperability.\nAgents discover and communicate with each other.\nWorks across organizations and platforms."),
]
for i, (abbr, full, desc) in enumerate(protocols):
    x = Inches(0.5 + i * 4.2)
    y = Inches(1.7)
    bg = add_bg_rect(slide, LIGHT_BLUE, x, y, Inches(3.9), Inches(4.8))
    bg.shadow.inherit = False
    add_text_box(slide, x + Inches(0.3), y + Inches(0.3), Inches(3.4), Inches(0.7),
                 abbr, font_size=36, bold=True, color=MSFT_BLUE, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.3), y + Inches(1.1), Inches(3.4), Inches(0.5),
                 full, font_size=15, color=MED_GRAY, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.3), y + Inches(1.8), Inches(3.4), Inches(2.5),
                 desc, font_size=15, color=DARK_GRAY)
add_notes(slide, "No vendor lock-in. These protocols work across providers. MCP is becoming the industry standard for tool integration.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10: Poll
# ══════════════════════════════════════════════════════════════════════════════
slide = content_slide("Where Are You on the AI Agent Journey?")
add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(0.6),
             "Quick poll — raise your hand or use the chat:", font_size=20, color=MED_GRAY)
options = [
    ("🟢", "I've built chatbots or copilots before"),
    ("🟡", "I've experimented with AI APIs (OpenAI, Azure, etc.)"),
    ("🔴", "This is my first time working with AI agents"),
]
for i, (dot, text) in enumerate(options):
    y = Inches(2.5 + i * 1.5)
    bg = add_bg_rect(slide, LIGHT_BLUE, Inches(1.5), y, Inches(10), Inches(1.1))
    bg.shadow.inherit = False
    add_text_box(slide, Inches(1.8), y + Inches(0.2), Inches(0.6), Inches(0.6),
                 dot, font_size=32)
    add_text_box(slide, Inches(2.6), y + Inches(0.2), Inches(8.5), Inches(0.6),
                 text, font_size=22, color=DARK_BLUE)
add_notes(slide, "Use Mentimeter, Slido, or show of hands. This helps you gauge the room and adjust depth accordingly.")


# ══════════════════════════════════════════════════════════════════════════════
# MODULE 1 SECTION DIVIDER
# ══════════════════════════════════════════════════════════════════════════════
section_title_slide("Module 1", "Agent Fundamentals", "09:30 – 10:15  ·  Theory + Lab")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11: AIAgent Abstraction
# ══════════════════════════════════════════════════════════════════════════════
slide = content_slide("The AIAgent Abstraction")
items = [
    "An Agent = LLM + Instructions + Tools",
    "AIAgent is the core type in the framework",
    "RunAsync() — get the complete response at once",
    "RunStreamingAsync() — receive tokens as they're generated",
    "Provider-agnostic: same code works with Azure OpenAI, Ollama, Anthropic…",
]
add_bullet_slide(slide, items, Inches(0.6), Inches(1.5), Inches(6), Inches(5),
                 font_size=19, color=DARK_GRAY)
# Diagram on right
bg = add_bg_rect(slide, LIGHT_BLUE, Inches(7.5), Inches(1.7), Inches(5.2), Inches(4.5))
bg.shadow.inherit = False
add_text_box(slide, Inches(7.8), Inches(2.0), Inches(4.6), Inches(4),
             "        ┌──────────────┐\n"
             "        │   AIAgent    │\n"
             "        ├──────────────┤\n"
             "        │ Instructions │\n"
             "        │ Provider     │\n"
             "        │ Tools [ ]    │\n"
             "        └──────┬───────┘\n"
             "               │\n"
             "        RunAsync()\n"
             "        RunStreamingAsync()",
             font_size=14, color=DARK_BLUE, font_name="Cascadia Code")
add_notes(slide, "Emphasize simplicity — you can build an agent in ~10 lines of code. The abstraction hides the complexity of LLM interaction.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12: First Agent Code
# ══════════════════════════════════════════════════════════════════════════════
slide = content_slide("10 Lines to Your First Agent")
code = '''AIAgent agent = new AzureOpenAIClient(
        new Uri(endpoint), new AzureCliCredential())
    .GetChatClient("gpt-4o-mini")
    .AsAIAgent(
        instructions: "You are a friendly assistant.",
        name: "HelloAgent");

// Complete response
Console.WriteLine(await agent.RunAsync("What is the largest city in France?"));

// Or stream token-by-token
await foreach (var update in agent.RunStreamingAsync("Tell me a fun fact."))
    Console.Write(update);'''
add_code_block(slide, code, Inches(0.6), Inches(1.5), Inches(12), Inches(5.3))
add_notes(slide, "Live-code this! Show it running. Point out: instructions shape behavior, provider is configured once, RunAsync vs RunStreamingAsync.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13: Provider Flexibility
# ══════════════════════════════════════════════════════════════════════════════
slide = content_slide("Swap Models, Keep Your Code")
table_slide_helper(slide,
    ["Provider", "Status", "Notes"],
    [
        ["Azure OpenAI", "✅ Fully supported", "Chat Completions + Responses API"],
        ["OpenAI", "✅ Fully supported", "Direct OpenAI API"],
        ["Azure AI Foundry", "✅ Fully supported", "Managed agent service"],
        ["Anthropic", "✅ Supported", "Claude models"],
        ["Ollama", "✅ Supported", "Local / on-premises models"],
        ["GitHub Copilot", "✅ Supported", "Copilot as a provider"],
    ],
    Inches(1.0), Inches(1.7), [3.5, 3.0, 5.0])
add_notes(slide, "The same agent code works with any provider. Demo swapping Azure OpenAI for Ollama if time permits. This is a key differentiator.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14: Streaming vs Non-Streaming
# ══════════════════════════════════════════════════════════════════════════════
slide = content_slide("Two Ways to Get Responses")
# Left box
bg = add_bg_rect(slide, LIGHT_BLUE, Inches(0.5), Inches(1.7), Inches(5.8), Inches(5.0))
bg.shadow.inherit = False
add_text_box(slide, Inches(0.8), Inches(1.9), Inches(5.2), Inches(0.6),
             "RunAsync()", font_size=24, bold=True, color=DARK_BLUE, font_name="Cascadia Code")
add_text_box(slide, Inches(0.8), Inches(2.7), Inches(5.2), Inches(3.5),
             "• Waits for the complete response\n"
             "• Returns the full text at once\n"
             "• Simpler to work with\n"
             "• Best for: background tasks,\n  non-interactive use cases",
             font_size=17, color=DARK_GRAY)

# Right box
bg = add_bg_rect(slide, RGBColor(0xE8, 0xF5, 0xE9), Inches(7.0), Inches(1.7), Inches(5.8), Inches(5.0))
bg.shadow.inherit = False
add_text_box(slide, Inches(7.3), Inches(1.9), Inches(5.2), Inches(0.6),
             "RunStreamingAsync()", font_size=24, bold=True, color=ACCENT_GREEN, font_name="Cascadia Code")
add_text_box(slide, Inches(7.3), Inches(2.7), Inches(5.2), Inches(3.5),
             "• Yields tokens as they're generated\n"
             "• Feels much faster to users\n"
             "• Uses IAsyncEnumerable<T>\n"
             "• Best for: chat UIs, real-time\n  interactive experiences\n\n"
             "⭐ Preferred for user-facing apps",
             font_size=17, color=DARK_GRAY)
add_notes(slide, "Streaming is almost always preferred for user-facing apps. The perceived latency is dramatically lower.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15: Lab 1
# ══════════════════════════════════════════════════════════════════════════════
slide = content_slide("🔬  Lab 1: Hello Agent")
add_text_box(slide, Inches(0.6), Inches(1.5), Inches(11), Inches(0.6),
             "Build your first agent in 15 minutes", font_size=22, bold=True, color=DARK_BLUE)
items = [
    "1.  Create a new console project and add MAF packages",
    "2.  Write an agent with custom instructions",
    "3.  Run it — non-streaming and streaming",
    "4.  Experiment with different personas (pirate, poet, expert…)",
    "5.  (Stretch) Build an interactive console chat loop",
]
add_bullet_slide(slide, items, Inches(0.6), Inches(2.5), Inches(8), Inches(4),
                 font_size=20)

bg = add_bg_rect(slide, LIGHT_BLUE, Inches(9.0), Inches(2.5), Inches(3.8), Inches(3.5))
bg.shadow.inherit = False
add_text_box(slide, Inches(9.2), Inches(2.7), Inches(3.4), Inches(3.2),
             "📁 Lab folder:\nlabs/lab1-hello-agent/\n\n⏱️ 15 minutes\n\n✅ Success:\nAgent responds &\nstreams tokens",
             font_size=15, color=DARK_BLUE)
add_notes(slide, "15 minutes. Circulate and help with Azure credential issues. Have az login ready as a troubleshooting step.")


# ══════════════════════════════════════════════════════════════════════════════
# MODULE 2 SECTION DIVIDER
# ══════════════════════════════════════════════════════════════════════════════
section_title_slide("Module 2", "Tools & Function Calling", "10:30 – 11:15  ·  Theory + Lab", ACCENT_GREEN)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16: Why Tools?
# ══════════════════════════════════════════════════════════════════════════════
slide = content_slide("Agents Without Tools Are Just Chatbots")
add_text_box(slide, Inches(0.6), Inches(1.6), Inches(12), Inches(0.7),
             "Tools give agents the ability to DO things — not just TALK about them.",
             font_size=22, color=DARK_BLUE)
# Two columns: without tools vs with tools
bg1 = add_bg_rect(slide, RGBColor(0xFF, 0xEB, 0xEE), Inches(0.5), Inches(2.8), Inches(5.8), Inches(3.8))
bg1.shadow.inherit = False
add_text_box(slide, Inches(0.8), Inches(3.0), Inches(5.2), Inches(0.5),
             "❌  Without Tools", font_size=20, bold=True, color=ACCENT_RED)
add_text_box(slide, Inches(0.8), Inches(3.7), Inches(5.2), Inches(2.5),
             '"What\'s the weather in Amsterdam?"\n\n→ "I don\'t have access to real-time\n   weather data, but typically..."',
             font_size=17, color=DARK_GRAY)

bg2 = add_bg_rect(slide, RGBColor(0xE8, 0xF5, 0xE9), Inches(7.0), Inches(2.8), Inches(5.8), Inches(3.8))
bg2.shadow.inherit = False
add_text_box(slide, Inches(7.3), Inches(3.0), Inches(5.2), Inches(0.5),
             "✅  With Tools", font_size=20, bold=True, color=ACCENT_GREEN)
add_text_box(slide, Inches(7.3), Inches(3.7), Inches(5.2), Inches(2.5),
             '"What\'s the weather in Amsterdam?"\n\n→ Calls GetWeather("Amsterdam")\n→ "It\'s 15°C and cloudy in\n   Amsterdam right now."',
             font_size=17, color=DARK_GRAY)
add_notes(slide, "This is the 'aha moment' module. When participants see the LLM autonomously decide to call their function, it clicks.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 17: Function Tools Code
# ══════════════════════════════════════════════════════════════════════════════
slide = content_slide("Your Code → Agent's Superpower")
code = '''// 1. Define a tool — just a method with [Description] attributes
[Description("Get the current weather for a given location.")]
static string GetWeather(
    [Description("The city name, e.g. 'Amsterdam'.")] string location)
    => $"The weather in {location} is cloudy, 15°C.";

// 2. Register the tool with the agent
AIAgent agent = new AzureOpenAIClient(new Uri(endpoint), new AzureCliCredential())
    .GetChatClient("gpt-4o-mini")
    .AsAIAgent(
        instructions: "You are a helpful weather assistant.",
        tools: [AIFunctionFactory.Create(GetWeather)]);

// 3. The agent calls the tool automatically when needed!
Console.WriteLine(await agent.RunAsync("What's the weather in Amsterdam?"));'''
add_code_block(slide, code, Inches(0.6), Inches(1.5), Inches(12), Inches(5.3))
add_notes(slide, "Stress: descriptions are CRITICAL. The LLM reads them to decide when and how to call the tool. Bad descriptions = bad tool usage.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 18: How Function Calling Works
# ══════════════════════════════════════════════════════════════════════════════
slide = content_slide("The LLM Decides — You Don't Hard-Code")
steps = [
    ('1', 'User asks a question', '"What\'s the weather in Amsterdam?"', MSFT_BLUE),
    ('2', 'LLM inspects available tools', 'Sees GetWeather(location) — matches!', ACCENT_PURPLE),
    ('3', 'Framework executes the function', 'GetWeather("Amsterdam") → "cloudy, 15°C"', ACCENT_GREEN),
    ('4', 'LLM crafts the response', '"It\'s currently 15°C and cloudy in Amsterdam."', MSFT_BLUE),
]
for i, (num, title, desc, clr) in enumerate(steps):
    y = Inches(1.6 + i * 1.4)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), y, Inches(0.65), Inches(0.65))
    circle.fill.solid()
    circle.fill.fore_color.rgb = clr
    circle.line.fill.background()
    add_text_box(slide, Inches(0.8), y + Inches(0.1), Inches(0.65), Inches(0.5),
                 num, font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.8), y + Inches(0.05), Inches(4), Inches(0.5),
                 title, font_size=19, bold=True, color=DARK_BLUE)
    add_text_box(slide, Inches(6.5), y + Inches(0.05), Inches(6), Inches(0.5),
                 desc, font_size=17, color=MED_GRAY)
add_notes(slide, "The magic: the LLM autonomously decided to call the function. You never wrote 'if user asks about weather, call GetWeather'.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 19: Tool Types
# ══════════════════════════════════════════════════════════════════════════════
slide = content_slide("The MAF Tool Ecosystem")
table_slide_helper(slide,
    ["Tool Type", "Description", "Provider Support"],
    [
        ["Function Tools", "Your own C#/Python methods", "All providers ✅"],
        ["Code Interpreter", "Sandboxed code execution", "Responses, Assistants, Foundry"],
        ["File Search", "Search uploaded documents", "Responses, Assistants, Foundry"],
        ["Web Search", "Live web search", "Chat Completion, Responses"],
        ["Hosted MCP Tools", "Microsoft-hosted MCP servers", "Responses, Foundry, Anthropic"],
        ["Local MCP Tools", "Your own / 3rd-party MCP servers", "All providers with functions ✅"],
    ],
    Inches(0.8), Inches(1.7), [3.0, 4.5, 4.0])
add_notes(slide, "Function tools are 80% of what you'll use day-to-day. MCP is the future — standardized tool integration across the industry.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 20: Agent-as-Tool
# ══════════════════════════════════════════════════════════════════════════════
slide = content_slide("Agents Calling Agents", "The Agent-as-a-Tool Pattern")
code = '''// Specialist agent
AIAgent weatherAgent = client.AsAIAgent(
    instructions: "You answer weather questions.",
    name: "WeatherAgent",
    description: "Answers weather questions.",
    tools: [AIFunctionFactory.Create(GetWeather)]);

// Main agent delegates to specialists
AIAgent mainAgent = client.AsAIAgent(
    instructions: "You are a helpful assistant. Delegate as needed.",
    tools: [weatherAgent.AsAIFunction()]);  // Agent becomes a tool!

mainAgent.RunAsync("What's the weather in Paris?");
// → mainAgent calls weatherAgent → weatherAgent calls GetWeather'''
add_code_block(slide, code, Inches(0.6), Inches(1.5), Inches(12), Inches(5.3))
add_notes(slide, "This is the simplest multi-agent pattern. The main agent doesn't know HOW to get weather — it delegates to a specialist.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 21: MCP
# ══════════════════════════════════════════════════════════════════════════════
slide = content_slide("MCP — Model Context Protocol", "The universal adapter for agent tools")
items = [
    "MCP — standardized tool discovery & invocation",
    "Hosted MCP — Microsoft manages the server (via Azure AI Foundry)",
    "Local MCP — you run your own tool server (any language, any platform)",
    "Any MCP client works with any MCP server — no vendor lock-in",
    "Rapidly becoming the industry standard (Anthropic, OpenAI, Microsoft, Google…)",
]
add_bullet_slide(slide, items, Inches(0.6), Inches(1.6), Inches(12), Inches(3.5),
                 font_size=20, color=DARK_GRAY)
add_text_box(slide, Inches(0.6), Inches(5.5), Inches(12), Inches(0.5),
             "💡 Think of MCP as USB-C for AI tools — one standard plug, everything connects.",
             font_size=18, bold=True, color=MSFT_BLUE)
add_notes(slide, "MCP is becoming the industry standard. Brief mention today — not a lab topic, but important context.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 22: Lab 2
# ══════════════════════════════════════════════════════════════════════════════
slide = content_slide("🔬  Lab 2: Agent with Tools")
add_text_box(slide, Inches(0.6), Inches(1.5), Inches(11), Inches(0.6),
             "Give your agent superpowers — 25 minutes", font_size=22, bold=True, color=DARK_BLUE)
items = [
    "1.  Define a function tool with [Description] attributes",
    "2.  Register it with the agent",
    "3.  Ask a question that triggers the tool — observe the magic!",
    "4.  Add a second tool — see the agent use both",
    "5.  (Stretch) Try Agent-as-a-Tool: build specialist agents",
    "6.  (Stretch) Experiment with bad tool descriptions",
]
add_bullet_slide(slide, items, Inches(0.6), Inches(2.5), Inches(8), Inches(4), font_size=20)
bg = add_bg_rect(slide, LIGHT_BLUE, Inches(9.0), Inches(2.5), Inches(3.8), Inches(3.5))
bg.shadow.inherit = False
add_text_box(slide, Inches(9.2), Inches(2.7), Inches(3.4), Inches(3.2),
             "📁 Lab folder:\nlabs/lab2-tools/\n\n⏱️ 25 minutes\n\n✅ Success:\nAgent calls your\nfunction automatically",
             font_size=15, color=DARK_BLUE)
add_notes(slide, "25 minutes. Walk the room during the 'bad description' exercise — it's a great teaching moment.")


# ══════════════════════════════════════════════════════════════════════════════
# MODULE 3 SECTION DIVIDER
# ══════════════════════════════════════════════════════════════════════════════
section_title_slide("Module 3", "Conversations, Memory\n& Middleware", "11:15 – 12:00  ·  Theory + Lab", ACCENT_PURPLE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 23: Memory Problem
# ══════════════════════════════════════════════════════════════════════════════
slide = content_slide("Agents Forget (Unless You Help)")
# Without session
bg1 = add_bg_rect(slide, RGBColor(0xFF, 0xEB, 0xEE), Inches(0.5), Inches(1.7), Inches(5.8), Inches(5.0))
bg1.shadow.inherit = False
add_text_box(slide, Inches(0.8), Inches(1.9), Inches(5.2), Inches(0.5),
             "❌  Without Session", font_size=20, bold=True, color=ACCENT_RED)
add_text_box(slide, Inches(0.8), Inches(2.6), Inches(5.2), Inches(3.5),
             'agent.RunAsync("My name is Alice.")\n→ "Nice to meet you, Alice!"\n\n'
             'agent.RunAsync("What is my name?")\n→ "I don\'t know your name.\n   Could you tell me?"',
             font_size=16, color=DARK_GRAY, font_name="Cascadia Code")

# With session
bg2 = add_bg_rect(slide, RGBColor(0xE8, 0xF5, 0xE9), Inches(7.0), Inches(1.7), Inches(5.8), Inches(5.0))
bg2.shadow.inherit = False
add_text_box(slide, Inches(7.3), Inches(1.9), Inches(5.2), Inches(0.5),
             "✅  With Session", font_size=20, bold=True, color=ACCENT_GREEN)
add_text_box(slide, Inches(7.3), Inches(2.6), Inches(5.2), Inches(3.5),
             'var session = await agent\n    .CreateSessionAsync();\n\n'
             'agent.RunAsync("My name is Alice.",\n    session);\n→ "Nice to meet you, Alice!"\n\n'
             'agent.RunAsync("What is my name?",\n    session);\n→ "Your name is Alice!"',
             font_size=16, color=DARK_GRAY, font_name="Cascadia Code")
add_notes(slide, "Show the contrast live. This is a visceral demo that makes the concept stick.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 24: Session code
# ══════════════════════════════════════════════════════════════════════════════
slide = content_slide("Sessions Keep the Conversation Alive")
code = '''// Create a session — maintains conversation history
AgentSession session = await agent.CreateSessionAsync();

// Turn 1: introduce yourself
Console.WriteLine(await agent.RunAsync(
    "My name is Alice and I love hiking.", session));

// Turn 2: the agent remembers context
Console.WriteLine(await agent.RunAsync(
    "What do you remember about me?", session));

// Session serialization for persistence
var serialized = agent.SerializeSession(session);
AgentSession restored = await agent.DeserializeSessionAsync(serialized);'''
add_code_block(slide, code, Inches(0.6), Inches(1.5), Inches(12), Inches(5.3))
add_notes(slide, "Live-demo: 'My name is Alice' → 'What's my name?' with and without session. Session serialization enables persistence across restarts.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 25: Context Providers
# ══════════════════════════════════════════════════════════════════════════════
slide = content_slide("Beyond Chat History — Injecting Knowledge")
items = [
    "Context Providers — inject persistent information into every agent run",
    "Use cases: user preferences, organization rules, knowledge base entries",
    "ChatHistoryProvider — custom storage for conversation history",
    "InMemoryHistoryProvider (default) vs CosmosDB, SQL, Redis (production)",
    "Context Providers vs RAG: providers add context per-turn, RAG retrieves on-demand",
]
add_bullet_slide(slide, items, Inches(0.6), Inches(1.6), Inches(12), Inches(3),
                 font_size=19, color=DARK_GRAY, bold_prefix=True)
add_text_box(slide, Inches(0.6), Inches(5.3), Inches(12), Inches(0.5),
             "💡 Context providers are how you bring domain knowledge to agents — think of them as \"always-on memory\".",
             font_size=17, bold=True, color=MSFT_BLUE)
add_notes(slide, "Context providers are the enterprise way to manage agent memory. For the workshop, InMemory is fine; mention CosmosDB for production.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 26: Middleware Pipeline
# ══════════════════════════════════════════════════════════════════════════════
slide = content_slide("Middleware: Intercept Everything")
# Pipeline visualization
pipeline_items = [
    ("Request", MSFT_BLUE),
    ("→", DARK_GRAY),
    ("Security", ACCENT_RED),
    ("→", DARK_GRAY),
    ("Logging", ACCENT_PURPLE),
    ("→", DARK_GRAY),
    ("Agent", ACCENT_GREEN),
    ("→", DARK_GRAY),
    ("Logging", ACCENT_PURPLE),
    ("→", DARK_GRAY),
    ("Response", MSFT_BLUE),
]
x = Inches(0.3)
for label, clr in pipeline_items:
    w = Inches(0.4) if label == "→" else Inches(1.1)
    if label != "→":
        bg = add_bg_rect(slide, clr, x, Inches(2.0), w, Inches(0.65))
        bg.shadow.inherit = False
        add_text_box(slide, x, Inches(2.05), w, Inches(0.55),
                     label, font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    else:
        add_text_box(slide, x, Inches(2.0), w, Inches(0.55),
                     "→", font_size=20, bold=True, color=DARK_GRAY, align=PP_ALIGN.CENTER)
    x += w

# Three types
add_text_box(slide, Inches(0.6), Inches(3.3), Inches(12), Inches(0.5),
             "Three types of middleware:", font_size=20, bold=True, color=DARK_BLUE)
types = [
    ("Agent Run Middleware", "Intercept input/output of every agent.RunAsync() call"),
    ("Function Calling Middleware", "Intercept every tool invocation — log, modify, or block"),
    ("IChatClient Middleware", "Intercept raw LLM calls for token-level control"),
]
for i, (title, desc) in enumerate(types):
    y = Inches(4.1 + i * 1.0)
    add_text_box(slide, Inches(0.8), y, Inches(5), Inches(0.5),
                 title, font_size=18, bold=True, color=DARK_BLUE)
    add_text_box(slide, Inches(6.0), y, Inches(7), Inches(0.5),
                 desc, font_size=16, color=MED_GRAY)
add_notes(slide, "Three types: Agent Run, Function Calling, IChatClient. They form a chain — order matters (security before logging!).")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 27: Middleware Use Cases
# ══════════════════════════════════════════════════════════════════════════════
slide = content_slide("What Can Middleware Do?")
use_cases = [
    ("🔒", "Security", "Block sensitive queries, validate input"),
    ("📊", "Telemetry", "Log latency, token counts, usage metrics"),
    ("🛡️", "Content Filtering", "Check responses for policy violations"),
    ("🔄", "Retry Logic", "Retry failed LLM calls with backoff"),
    ("📝", "Audit", "Log all tool calls for compliance"),
    ("⏱️", "Rate Limiting", "Throttle requests per user/session"),
]
for i, (icon, title, desc) in enumerate(use_cases):
    col = i % 3
    row = i // 3
    x = Inches(0.4 + col * 4.2)
    y = Inches(1.7 + row * 2.7)
    bg = add_bg_rect(slide, LIGHT_BLUE, x, y, Inches(3.9), Inches(2.2))
    bg.shadow.inherit = False
    add_text_box(slide, x + Inches(0.2), y + Inches(0.2), Inches(0.6), Inches(0.5),
                 icon, font_size=28)
    add_text_box(slide, x + Inches(0.9), y + Inches(0.2), Inches(2.8), Inches(0.5),
                 title, font_size=19, bold=True, color=DARK_BLUE)
    add_text_box(slide, x + Inches(0.2), y + Inches(1.0), Inches(3.5), Inches(1.0),
                 desc, font_size=16, color=MED_GRAY)
add_notes(slide, "Security teams love middleware — it's how you make agents production-safe. Emphasize that middleware is the enterprise plumbing layer.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 28: Lab 3
# ══════════════════════════════════════════════════════════════════════════════
slide = content_slide("🔬  Lab 3: Multi-Turn Agent with Middleware")
add_text_box(slide, Inches(0.6), Inches(1.5), Inches(11), Inches(0.6),
             "Stateful conversations + enterprise plumbing — 25 minutes", font_size=22, bold=True, color=DARK_BLUE)
items = [
    "1.  Create a session and have a multi-turn conversation",
    "2.  Verify the agent remembers context across turns",
    "3.  Add logging middleware (message counts, turn tracking)",
    "4.  Add function calling middleware (tool name, timing)",
    "5.  (Stretch) Build security middleware that blocks sensitive queries",
]
add_bullet_slide(slide, items, Inches(0.6), Inches(2.5), Inches(8), Inches(4), font_size=20)
bg = add_bg_rect(slide, LIGHT_BLUE, Inches(9.0), Inches(2.5), Inches(3.8), Inches(3.5))
bg.shadow.inherit = False
add_text_box(slide, Inches(9.2), Inches(2.7), Inches(3.4), Inches(3.2),
             "📁 Lab folder:\nlabs/lab3-conversations\n  -middleware/\n\n⏱️ 25 minutes\n\n✅ Success:\nMulti-turn works &\nmiddleware logs",
             font_size=15, color=DARK_BLUE)
add_notes(slide, "25 minutes. The security middleware exercise is great for discussion — how would you protect a production agent?")


# ══════════════════════════════════════════════════════════════════════════════
# MODULE 4 SECTION DIVIDER
# ══════════════════════════════════════════════════════════════════════════════
section_title_slide("Module 4", "Workflows & Multi-Agent\nOrchestration", "12:45 – 13:30  ·  Theory + Lab", ACCENT_ORANGE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 29: Agents vs Workflows
# ══════════════════════════════════════════════════════════════════════════════
slide = content_slide("When a Single Agent Isn't Enough")
table_slide_helper(slide,
    ["Use an Agent when…", "Use a Workflow when…"],
    [
        ["Task is open-ended or conversational", "Process has well-defined steps"],
        ["Need autonomous tool use & planning", "Need explicit control over execution order"],
        ["Single LLM call (+ tools) suffices", "Multiple agents/functions must coordinate"],
        ["Exploration & creativity needed", "Reliability & repeatability required"],
    ],
    Inches(0.8), Inches(1.7), [5.8, 5.8])
add_text_box(slide, Inches(0.8), Inches(5.5), Inches(12), Inches(0.5),
             '💡 Rule of thumb: "If you can draw a flowchart, use a workflow."',
             font_size=18, bold=True, color=MSFT_BLUE)
add_notes(slide, "If you can draw a flowchart, use a workflow. If the task is open-ended, use an agent. Most real apps use both.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 30: Workflow Building Blocks
# ══════════════════════════════════════════════════════════════════════════════
slide = content_slide("Executors + Edges = Graph")
items = [
    "Executors — the nodes (steps) in your workflow:",
    "    • Function Executor — runs your code",
    "    • Agent Executor — runs an LLM agent",
    "    • Nested Workflow — embeds a sub-workflow",
    "",
    "Edges — connect executors (arrows in the graph):",
    "    • Direct edges — always follow this path",
    "    • Conditional edges — choose path based on output",
    "",
    "WorkflowContext — typed message passing between nodes",
    "    • SendMessageAsync() — forward to next executor",
    "    • YieldOutputAsync() — final workflow output",
]
add_bullet_slide(slide, items, Inches(0.6), Inches(1.6), Inches(12), Inches(5.5),
                 font_size=18, color=DARK_GRAY)
add_notes(slide, "Type-safe messages flow along edges. The framework validates message types at build time, catching errors early.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 31: Workflow Patterns
# ══════════════════════════════════════════════════════════════════════════════
slide = content_slide("Common Multi-Agent Patterns")
patterns = [
    ("Sequential", "A → B → C", "Simple pipeline\nEach step feeds the next", MSFT_BLUE),
    ("Parallel", "A → [B, C] → D", "Fan-out / fan-in\nMultiple agents in parallel", ACCENT_GREEN),
    ("Conditional", "A → if X: B\n     else: C", "Route based on output\nDynamic execution paths", ACCENT_ORANGE),
    ("Supervisor", "Manager →\n  [B, C, D]", "One agent delegates\nto specialists", ACCENT_PURPLE),
]
for i, (name, diagram, desc, clr) in enumerate(patterns):
    col = i % 2
    row = i // 2
    x = Inches(0.4 + col * 6.5)
    y = Inches(1.6 + row * 2.8)
    bg = add_bg_rect(slide, clr, x, y, Inches(6.0), Inches(2.4))
    bg.shadow.inherit = False
    add_text_box(slide, x + Inches(0.3), y + Inches(0.2), Inches(2.5), Inches(0.5),
                 name, font_size=20, bold=True, color=WHITE)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.8), Inches(2.5), Inches(1.4),
                 diagram, font_size=15, color=WHITE, font_name="Cascadia Code")
    add_text_box(slide, x + Inches(3.2), y + Inches(0.5), Inches(2.5), Inches(1.6),
                 desc, font_size=15, color=WHITE)
add_notes(slide, "Most real workflows combine these patterns. Show diagram of a real-world example if possible.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 32: Checkpointing
# ══════════════════════════════════════════════════════════════════════════════
slide = content_slide("Long-Running Workflows Need State")
items = [
    "Checkpointing — save/restore workflow state at any point",
    "Human-in-the-Loop — pause workflow, get human approval, resume",
    "Durable execution — survive process restarts and failures",
    "Azure Functions integration — serverless, durable workflows",
]
add_bullet_slide(slide, items, Inches(0.6), Inches(1.6), Inches(12), Inches(3),
                 font_size=20, color=DARK_GRAY, bold_prefix=True)

# Visual: workflow with checkpoint
add_text_box(slide, Inches(1.5), Inches(4.5), Inches(10), Inches(2),
             "Input → [Agent A] → 💾 Checkpoint → ⏸️ Human Approval → [Agent B] → 💾 Checkpoint → Output",
             font_size=18, color=MSFT_BLUE, bold=True, align=PP_ALIGN.CENTER,
             font_name="Cascadia Code")
add_notes(slide, "Critical for production workflows that take hours/days. Think expense approvals, contract reviews, multi-stage processing.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 33: Hosting
# ══════════════════════════════════════════════════════════════════════════════
slide = content_slide("From Console to Production")
table_slide_helper(slide,
    ["Hosting Option", "Best For", "Key Feature"],
    [
        ["ASP.NET Core", "Web APIs, custom backends", "Full control, DI integration"],
        ["Azure Functions (Durable)", "Serverless, long-running", "Pay-per-use, auto-scale"],
        ["A2A Protocol", "Multi-agent systems", "Cross-system agent discovery"],
        ["AG-UI Protocol", "Web frontends", "Rich streaming UI experience"],
        ["OpenAI-Compatible Endpoints", "OpenAI-compatible clients", "Drop-in replacement"],
    ],
    Inches(0.8), Inches(1.7), [3.5, 3.5, 4.5])
add_notes(slide, "Brief overview — full hosting is beyond today's scope. Show the 5-line ASP.NET Core hosting code if time allows.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 34: Lab 4
# ══════════════════════════════════════════════════════════════════════════════
slide = content_slide("🔬  Lab 4: Build a Workflow")
add_text_box(slide, Inches(0.6), Inches(1.5), Inches(11), Inches(0.6),
             "Chain agents into an orchestrated pipeline — 25 minutes", font_size=22, bold=True, color=DARK_BLUE)
items = [
    "1.  Build a simple function workflow: UpperCase → Reverse",
    "2.  Observe data flowing between executors",
    "3.  Build an agent workflow: Research → Writer",
    "4.  See how two LLM agents collaborate via the graph",
    "5.  (Stretch) Add a Translator step to the pipeline",
    "6.  (Stretch) Add conditional routing or checkpointing",
]
add_bullet_slide(slide, items, Inches(0.6), Inches(2.5), Inches(8), Inches(4), font_size=20)
bg = add_bg_rect(slide, LIGHT_BLUE, Inches(9.0), Inches(2.5), Inches(3.8), Inches(3.5))
bg.shadow.inherit = False
add_text_box(slide, Inches(9.2), Inches(2.7), Inches(3.4), Inches(3.2),
             "📁 Lab folder:\nlabs/lab4-workflows/\n\n⏱️ 25 minutes\n\n✅ Success:\nData flows through\na multi-step graph",
             font_size=15, color=DARK_BLUE)
add_notes(slide, "If short on time, do the function workflow as guided walkthrough. The agent workflow is the impressive payoff.")


# ══════════════════════════════════════════════════════════════════════════════
# HACKATHON SECTION DIVIDER
# ══════════════════════════════════════════════════════════════════════════════
section_title_slide("Part 2", "🏗️  Hackathon\nTransform the Existing App", "13:30 – 17:00", ACCENT_ORANGE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 36: Existing App (placeholder)
# ══════════════════════════════════════════════════════════════════════════════
slide = content_slide("Meet [Existing App Name]", "CUSTOMIZE: Add screenshots and demo the existing application")
add_text_box(slide, Inches(2), Inches(3), Inches(9), Inches(2),
             "[ Insert existing app screenshots here ]\n\n"
             "Walk through the current functionality, user journeys,\n"
             "and pain points that agents could address.",
             font_size=20, color=MED_GRAY, align=PP_ALIGN.CENTER)
add_notes(slide, "CUSTOMIZE THIS SLIDE. Walk through the existing app live. Show the user journeys. Identify where agents add value.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 37: Challenge Scenarios
# ══════════════════════════════════════════════════════════════════════════════
slide = content_slide("The Challenge: Add an Agentic Experience")
scenarios = [
    ("💬", "Conversational\nAssistant", "Chat-based agent that helps\nusers navigate the app", "Agent, Tools, Session"),
    ("⚙️", "Process\nAutomation", "Automate a multi-step\nbusiness workflow", "Workflows, Human-\nin-the-Loop"),
    ("🔍", "Smart Search\n& Insights", "Agent that searches data\nand provides summaries", "Tools, Context\nProviders"),
    ("🤖", "Multi-Agent\nCollaboration", "Specialist agents working\ntogether on complex tasks", "Agent-as-Tool,\nWorkflows"),
]
for i, (icon, title, desc, concepts) in enumerate(scenarios):
    x = Inches(0.3 + i * 3.2)
    bg = add_bg_rect(slide, LIGHT_BLUE, x, Inches(1.6), Inches(3.0), Inches(5.3))
    bg.shadow.inherit = False
    add_text_box(slide, x + Inches(0.2), Inches(1.8), Inches(2.6), Inches(0.5),
                 icon, font_size=36, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), Inches(2.4), Inches(2.6), Inches(0.8),
                 title, font_size=17, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), Inches(3.4), Inches(2.6), Inches(1.5),
                 desc, font_size=14, color=MED_GRAY, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), Inches(5.2), Inches(2.6), Inches(1.2),
                 f"MAF Concepts:\n{concepts}", font_size=12, color=MSFT_BLUE, align=PP_ALIGN.CENTER)
add_notes(slide, "Teams choose 1-2 scenarios. All are valid. Customize these scenarios to match the partner's application.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 38: Hackathon Rules
# ══════════════════════════════════════════════════════════════════════════════
slide = content_slide("Hackathon Ground Rules")
items = [
    "👥  Teams of 3–5 people",
    "⏱️  2.5 hours of hacking time (14:00 – 16:30)",
    "📦  Starter code & scaffolding provided (see hackathon/starter-guide-csharp.md or starter-guide-python.md)",
    "🧑‍🏫  Mentors circulating — ask for help!",
    "🎤  5-minute demo per team at the end",
    "",
    "Judging criteria:",
    "   ⭐ Creativity & usefulness of the agentic experience",
    "   ⭐ Correct use of MAF concepts (tools, sessions, workflows)",
    "   ⭐ Working demo (doesn't need to be polished!)",
    "   ⭐ Team collaboration & presentation",
]
add_bullet_slide(slide, items, Inches(0.6), Inches(1.6), Inches(12), Inches(5.5),
                 font_size=20, color=DARK_GRAY)
add_notes(slide, "Form teams now, pick scenarios, set up environments during the break. Emphasize: a working demo of 1 feature beats a broken demo of 3.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 39: Show & Tell
# ══════════════════════════════════════════════════════════════════════════════
slide = new_slide()
add_bg_rect(slide, DARK_BLUE)
add_text_box(slide, Inches(1), Inches(2.0), Inches(11), Inches(1.5),
             "🎤  Show & Tell", font_size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(4.0), Inches(11), Inches(1),
             "5 minutes per team  ·  Show your scenario  ·  Demo the agent  ·  Share one insight",
             font_size=22, color=RGBColor(0xBB, 0xDE, 0xFB), align=PP_ALIGN.CENTER)
add_notes(slide, "5 min per team. Keep it moving. Celebrate creativity and effort, not just polish.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 40: Key Takeaways
# ══════════════════════════════════════════════════════════════════════════════
slide = content_slide("What We Learned Today")
takeaways = [
    "Agents = LLM + Instructions + Tools + Memory",
    "Tools are the superpower — function calling is where the magic happens",
    "Sessions make agents stateful; middleware makes them production-ready",
    "Workflows give explicit control over multi-agent orchestration",
    "MAF is the unified successor to Semantic Kernel + AutoGen",
]
for i, text in enumerate(takeaways):
    y = Inches(1.7 + i * 1.05)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), y, Inches(0.55), Inches(0.55))
    circle.fill.solid()
    circle.fill.fore_color.rgb = MSFT_BLUE
    circle.line.fill.background()
    add_text_box(slide, Inches(0.7), y + Inches(0.08), Inches(0.55), Inches(0.45),
                 str(i + 1), font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.5), y + Inches(0.08), Inches(11), Inches(0.5),
                 text, font_size=20, color=DARK_GRAY)
add_notes(slide, "Reinforce the 'so what' — these aren't toys, they're production-grade tools. MAF is the path forward for enterprise agentic AI.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 41: Resources
# ══════════════════════════════════════════════════════════════════════════════
slide = content_slide("Continue Your Journey")
resources = [
    ("📚", "Official Docs", "learn.microsoft.com/agent-framework"),
    ("💻", "GitHub Repo", "github.com/microsoft/agent-framework"),
    ("🧪", "Step-by-Step Workshop", "github.com/warnov/ms-agent-framework-step-by-step-workshop"),
    ("🌐", "Agent Framework", "aka.ms/agentframework"),
    ("🔄", "Migration Guides", "Semantic Kernel → MAF  |  AutoGen → MAF"),
]
for i, (icon, title, url) in enumerate(resources):
    y = Inches(1.6 + i * 1.1)
    add_text_box(slide, Inches(0.7), y, Inches(0.6), Inches(0.5), icon, font_size=24)
    add_text_box(slide, Inches(1.5), y, Inches(3.5), Inches(0.5),
                 title, font_size=19, bold=True, color=DARK_BLUE)
    add_text_box(slide, Inches(5.2), y, Inches(7.5), Inches(0.5),
                 url, font_size=17, color=MSFT_BLUE, font_name="Cascadia Code")
add_notes(slide, "Send follow-up email with all links + today's code within 24 hours.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 42: Thank You
# ══════════════════════════════════════════════════════════════════════════════
slide = new_slide()
add_bg_rect(slide, DARK_BLUE)
add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
             "Thank You!", font_size=60, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(3.5), Inches(11), Inches(1),
             "[ Presenter Names  ·  Contact Info ]",
             font_size=22, color=RGBColor(0xBB, 0xDE, 0xFB), align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(5.0), Inches(11), Inches(1),
             "📋  Please fill out the feedback survey!\n[ QR Code / Link Here ]",
             font_size=20, color=RGBColor(0x99, 0xBB, 0xDD), align=PP_ALIGN.CENTER)
add_notes(slide, "Encourage honest feedback. Thank the partner. Remind about the follow-up email with materials.")


# ══════════════════════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════════════════════
output_path = r"C:\project\MAF-training\MAF-Workshop-Deck.pptx"
prs.save(output_path)
print(f"✅ Saved {len(prs.slides)} slides to {output_path}")
